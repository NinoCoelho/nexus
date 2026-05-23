#!/usr/bin/env python3
"""Generate the Nexus Pre-Seed Investor Pitch Deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
import math

# ---------------------------------------------------------------------------
# Design system
# ---------------------------------------------------------------------------
BG_DARK = RGBColor(0x0B, 0x14, 0x26)
BG_CARD = RGBColor(0x11, 0x1E, 0x36)
BG_CARD_ALT = RGBColor(0x15, 0x23, 0x40)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_LIGHT = RGBColor(0x60, 0xA5, 0xFA)
ACCENT_DIM = RGBColor(0x1E, 0x40, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
GREEN = RGBColor(0x10, 0xB9, 0x81)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
RED_SOFT = RGBColor(0xF8, 0x71, 0x71)

FONT = "Calibri"
SW = Inches(13.333)
SH = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def tb(slide, left, top, width, height, text, size=14, color=WHITE,
       bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT, anchor=MSO_ANCHOR.TOP):
    """Add a text box and return it."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    try:
        tf.paragraphs[0].space_after = Pt(0)
        tf.paragraphs[0].space_before = Pt(0)
    except Exception:
        pass
    return txBox


def multi_tb(slide, left, top, width, height, lines, anchor=MSO_ANCHOR.TOP):
    """Add a text box with multiple styled lines.
    lines = [(text, size, color, bold, alignment), ...]
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, (text, size, color, bold, align) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = FONT
        p.alignment = align
        p.space_after = Pt(4)
        p.space_before = Pt(2)
    return txBox


def card(slide, left, top, width, height, fill=BG_CARD, border_color=None):
    """Add a rounded rectangle card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def accent_line(slide, left, top, width, color=ACCENT, thickness=Pt(3)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, thickness
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def circle(slide, left, top, size, fill):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def slide_number(slide, num, total=16):
    tb(slide, Inches(12.3), Inches(7.05), Inches(0.9), Inches(0.35),
       f"{num} / {total}", size=9, color=GRAY, alignment=PP_ALIGN.RIGHT)


def section_title(slide, title, subtitle=None):
    accent_line(slide, Inches(0.75), Inches(0.55), Inches(0.6))
    tb(slide, Inches(0.75), Inches(0.7), Inches(11), Inches(0.6),
       title, size=32, color=WHITE, bold=True)
    if subtitle:
        tb(slide, Inches(0.75), Inches(1.25), Inches(11), Inches(0.4),
           subtitle, size=14, color=GRAY)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    # Subtle gradient bar at top
    rect(slide, Inches(0), Inches(0), SW, Inches(0.06), ACCENT)

    # Company name
    tb(slide, Inches(0.75), Inches(1.8), Inches(11.8), Inches(1.2),
       "NEXUS", size=80, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

    # Accent line
    accent_line(slide, Inches(0.75), Inches(3.05), Inches(1.8))

    # Tagline
    tb(slide, Inches(0.75), Inches(3.35), Inches(11), Inches(0.7),
       "The Self-Evolving AI Agent", size=28, color=ACCENT_LIGHT, bold=False)

    # Description
    tb(slide, Inches(0.75), Inches(4.2), Inches(8), Inches(1.0),
       "A personal AI platform that learns, creates tools, and improves\nitself — while your data stays private.",
       size=16, color=GRAY)

    # Footer
    tb(slide, Inches(0.75), Inches(6.3), Inches(5), Inches(0.3),
       "Pre-Seed Investment Deck  •  2026  •  Confidential",
       size=11, color=GRAY)

    # Decorative circle (top right)
    circle(slide, Inches(10.5), Inches(0.8), Inches(2.2), ACCENT_DIM)
    tb(slide, Inches(10.6), Inches(1.3), Inches(2.0), Inches(1.2),
       "AI", size=48, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    slide_number(slide, 1)


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    section_title(slide, "The Problem", "Why current AI assistants fall short")

    problems = [
        ("Stateless & Forgetful", ACCENT,
         ["No persistent memory across sessions",
          "Can't build on past work or learn preferences",
          "Every conversation starts from zero"]),
        ("Privacy Black Box", GOLD,
         ["All data sent to third-party clouds",
          "No control over how data is stored or used",
          "Enterprise compliance is impossible"]),
        ("One-Size-Fits-All", PURPLE,
         ["No real customization beyond prompt tweaks",
          "Can't create workflows or automate tasks",
          "No self-improvement or skill building"]),
    ]

    card_w = Inches(3.6)
    card_h = Inches(4.2)
    gap = Inches(0.55)
    start_x = Inches(0.75)
    start_y = Inches(2.0)

    for i, (title, color, bullets) in enumerate(problems):
        x = start_x + i * (card_w + gap)
        c = card(slide, x, start_y, card_w, card_h, BG_CARD, border_color=RGBColor(0x1E, 0x29, 0x3B))

        # Color bar at top of card
        rect(slide, x, start_y, card_w, Inches(0.06), color)

        # Icon circle
        circle(slide, x + Inches(0.3), start_y + Inches(0.35), Inches(0.6), color)
        icon_labels = ["01", "02", "03"]
        tb(slide, x + Inches(0.3), start_y + Inches(0.38), Inches(0.6), Inches(0.5),
           icon_labels[i], size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Title
        tb(slide, x + Inches(0.3), start_y + Inches(1.15), Inches(3.0), Inches(0.5),
           title, size=18, color=WHITE, bold=True)

        # Bullets
        for j, bullet in enumerate(bullets):
            tb(slide, x + Inches(0.3), start_y + Inches(1.75) + j * Inches(0.65),
               Inches(3.0), Inches(0.6),
               f"→  {bullet}", size=12, color=LIGHT)

    slide_number(slide, 2)


def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    section_title(slide, "The Solution", "Nexus: The AI agent that evolves with you")

    # Hero statement
    hero_card = card(slide, Inches(0.75), Inches(2.0), Inches(11.8), Inches(1.3),
                     BG_CARD, border_color=ACCENT_DIM)
    tb(slide, Inches(1.2), Inches(2.15), Inches(10.8), Inches(1.0),
       "The only AI platform where the agent creates its own tools, refines its behavior, "
       "and builds knowledge over time — all while keeping your data on your machine.",
       size=17, color=WHITE, bold=False)

    # 4 pillars
    pillars = [
        ("Self-Evolving", "Dream engine creates &\nrefines skills autonomously", ACCENT),
        ("Local-First Privacy", "Data stays on your device.\nCloud features are optional.", GREEN),
        ("Model-Agnostic", "Works with any LLM.\nGateway optimizes cost & quality.", GOLD),
        ("Full Personal OS", "Vault, kanban, calendar,\nvoice, data, knowledge graph.", PURPLE),
    ]

    pw = Inches(2.7)
    ph = Inches(2.8)
    gap = Inches(0.4)
    sx = Inches(0.75)
    sy = Inches(3.7)

    for i, (title, desc, color) in enumerate(pillars):
        x = sx + i * (pw + gap)
        c = card(slide, x, sy, pw, ph, BG_CARD)

        # Top accent bar
        rect(slide, x, sy, pw, Inches(0.05), color)

        # Title
        tb(slide, x + Inches(0.25), sy + Inches(0.3), Inches(2.2), Inches(0.5),
           title, size=18, color=color, bold=True)

        # Description
        tb(slide, x + Inches(0.25), sy + Inches(0.9), Inches(2.2), Inches(1.5),
           desc, size=13, color=LIGHT)

    slide_number(slide, 3)


def slide_product(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    section_title(slide, "Product Overview", "A complete AI agent platform, not just a chatbot")

    features = [
        ("Agent Loop", "20+ tools, sub-agent spawning,\ncontext compaction, overflow guard", ACCENT),
        ("Vault & Knowledge", "Markdown files, FTS5 search,\nbacklink graph, DuckDB datatables", GREEN),
        ("Dream Engine", "4-phase idle processing:\nconsolidate → insight → refine → rehearse", PURPLE),
        ("Voice & TTS", "Local Piper TTS + Whisper STT,\nvoice acknowledgments, multi-language", CYAN),
        ("Kanban & Calendar", "Vault-native boards, iCal calendar,\nalarm notifications, auto-dispatch", GOLD),
        ("Tunnel & Sharing", "Cloudflare tunnel, secure access\ncodes, read-only session sharing", RED_SOFT),
    ]

    cols = 3
    rows = 2
    cw = Inches(3.7)
    ch = Inches(2.2)
    gx = Inches(0.5)
    gy = Inches(0.4)
    sx = Inches(0.75)
    sy = Inches(2.0)

    for i, (title, desc, color) in enumerate(features):
        r = i // cols
        c_idx = i % cols
        x = sx + c_idx * (cw + gx)
        y = sy + r * (ch + gy)

        ca = card(slide, x, y, cw, ch, BG_CARD)

        # Left accent bar
        rect(slide, x, y, Inches(0.05), ch, color)

        # Title
        tb(slide, x + Inches(0.25), y + Inches(0.2), Inches(3.2), Inches(0.4),
           title, size=16, color=color, bold=True)

        # Description
        tb(slide, x + Inches(0.25), y + Inches(0.7), Inches(3.2), Inches(1.3),
           desc, size=12, color=LIGHT)

    # Bottom callout
    tb(slide, Inches(0.75), Inches(6.85), Inches(11.8), Inches(0.35),
       "50+ bundled skills  •  9 UI views  •  Docker, macOS, Windows  •  Fully offline capable",
       size=12, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)

    slide_number(slide, 4)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    section_title(slide, "Architecture", "Model-agnostic by design, optimized by the Nexus Gateway")

    # Flow: boxes with arrows
    flow_items = [
        ("Local Client\n(Nexus App)", "Self-hosted, 127.0.0.1\nVault, Skills, Memory", ACCENT),
        ("Nexus Gateway\n(LiteLLM Proxy)", "Intelligent routing\nCost & quality optimization", PURPLE),
        ("Any LLM\n(GLM 5.1, GPT-4, ...)", "OpenAI, Anthropic, Google,\nDeepSeek, local models", GREEN),
    ]

    bw = Inches(3.0)
    bh = Inches(2.5)
    total_w = 3 * bw + 2 * Inches(1.2)
    sx = (SW - total_w) / 2
    sy = Inches(2.5)

    for i, (title, desc, color) in enumerate(flow_items):
        x = sx + i * (bw + Inches(1.2))
        c = card(slide, x, sy, bw, bh, BG_CARD, border_color=color)

        tb(slide, x + Inches(0.2), sy + Inches(0.25), bw - Inches(0.4), Inches(0.7),
           title, size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        tb(slide, x + Inches(0.2), sy + Inches(1.1), bw - Inches(0.4), Inches(1.2),
           desc, size=11, color=LIGHT, alignment=PP_ALIGN.CENTER)

        # Arrow between boxes
        if i < 2:
            arrow_x = x + bw + Inches(0.15)
            arrow_y = sy + bh / 2 - Inches(0.03)
            rect(slide, arrow_x, arrow_y, Inches(0.9), Inches(0.06), color)

    # Key insight box
    insight_card = card(slide, Inches(1.5), Inches(5.6), Inches(10.3), Inches(1.1),
                        BG_CARD_ALT, border_color=ACCENT_DIM)
    tb(slide, Inches(1.8), Inches(5.7), Inches(0.15), Inches(0.15),
       "", size=10, color=ACCENT)
    multi_tb(slide, Inches(2.0), Inches(5.7), Inches(9.5), Inches(0.9), [
        ("Cost Advantage:  GLM 5.1 backbone is 60-80% cheaper than GPT-4/Claude at retail.  "
         "LiteLLM routes each query to the optimal model.", 13, ACCENT_LIGHT, True, PP_ALIGN.LEFT),
        ("Users never think about models — Nexus picks the best one for the task, at the lowest cost.", 12, GRAY, False, PP_ALIGN.LEFT),
    ])

    slide_number(slide, 5)


def slide_dream(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    section_title(slide, "Self-Evolution: The Dream Engine",
                  "The only AI agent that autonomously improves itself during idle time")

    phases = [
        ("1", "Consolidate", "Deduplicate memories,\nfix dates, prune stale data", ACCENT),
        ("2", "Extract Insights", "Cross-session pattern\nrecognition & learning", CYAN),
        ("3", "Refine Skills", "Suggest improvements\nto existing skills", PURPLE),
        ("4", "Rehearse", "Test scenarios against\nrefined capabilities", GREEN),
    ]

    pw = Inches(2.5)
    ph = Inches(2.8)
    gap = Inches(0.4)
    total = 4 * pw + 3 * gap
    sx = (SW - total) / 2
    sy = Inches(2.2)

    for i, (num, title, desc, color) in enumerate(phases):
        x = sx + i * (pw + gap)
        c = card(slide, x, sy, pw, ph, BG_CARD)

        # Number circle
        circle(slide, x + pw / 2 - Inches(0.3), sy + Inches(0.25), Inches(0.6), color)
        tb(slide, x + pw / 2 - Inches(0.3), sy + Inches(0.28), Inches(0.6), Inches(0.5),
           num, size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Title
        tb(slide, x + Inches(0.15), sy + Inches(1.05), pw - Inches(0.3), Inches(0.4),
           title, size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        tb(slide, x + Inches(0.15), sy + Inches(1.6), pw - Inches(0.3), Inches(1.0),
           desc, size=11, color=LIGHT, alignment=PP_ALIGN.CENTER)

        # Connector arrow
        if i < 3:
            ax = x + pw + Inches(0.05)
            ay = sy + ph / 2
            rect(slide, ax, ay, gap - Inches(0.1), Inches(0.05), GRAY)

    # Bottom callout
    callout = card(slide, Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.2), BG_CARD_ALT)
    multi_tb(slide, Inches(1.3), Inches(5.6), Inches(10.7), Inches(1.0), [
        ("Why this matters:", 14, WHITE, True, PP_ALIGN.LEFT),
        ("Skills are staged for user review — the agent never auto-deploys.  "
         "Guard rails scan for credential exfiltration, destructive patterns, and prompt injection.  "
         "Trust tiers (builtin / user / agent) control permissions.", 12, LIGHT, False, PP_ALIGN.LEFT),
    ])

    slide_number(slide, 6)


def slide_market(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    section_title(slide, "Market Opportunity", "A massive and rapidly growing market")

    # TAM / SAM / SOM as concentric cards
    # Outer: TAM
    outer = card(slide, Inches(0.75), Inches(2.0), Inches(5.8), Inches(4.8), BG_CARD)
    tb(slide, Inches(1.0), Inches(2.15), Inches(5.3), Inches(0.4),
       "TAM — Total AI Agent Market", size=14, color=ACCENT_LIGHT, bold=True)
    tb(slide, Inches(1.0), Inches(2.55), Inches(5.3), Inches(0.6),
       "$100B+", size=44, color=WHITE, bold=True)
    tb(slide, Inches(1.0), Inches(3.2), Inches(5.3), Inches(0.8),
       "The global market for AI agents and autonomous\nassistants by 2028 (IDC, Grand View Research).",
       size=11, color=GRAY)

    # Middle: SAM
    mid = card(slide, Inches(1.3), Inches(4.1), Inches(4.7), Inches(2.4), BG_CARD_ALT)
    tb(slide, Inches(1.55), Inches(4.2), Inches(4.2), Inches(0.35),
       "SAM — Personal AI Platforms", size=13, color=CYAN, bold=True)
    tb(slide, Inches(1.55), Inches(4.55), Inches(4.2), Inches(0.5),
       "$12B", size=36, color=WHITE, bold=True)
    tb(slide, Inches(1.55), Inches(5.1), Inches(4.2), Inches(0.6),
       "Self-hosted and privacy-first AI agent platforms\nfor prosumers and SMBs.",
       size=11, color=GRAY)

    # Inner: SOM
    inner = card(slide, Inches(1.8), Inches(5.6), Inches(3.7), Inches(1.0), ACCENT_DIM)
    tb(slide, Inches(2.0), Inches(5.65), Inches(3.3), Inches(0.3),
       "SOM — Year 3 Reach", size=12, color=WHITE, bold=True)
    tb(slide, Inches(2.0), Inches(5.95), Inches(3.3), Inches(0.4),
       "$400M   (45K subscribers + enterprise)", size=14, color=WHITE, bold=True)

    # Right side: market drivers
    drivers_card = card(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.8), BG_CARD)
    tb(slide, Inches(7.3), Inches(2.15), Inches(5.0), Inches(0.4),
       "Market Tailwinds", size=16, color=WHITE, bold=True)

    drivers = [
        ("Agent-First Computing", "Shift from chatbots to autonomous agents doing real work"),
        ("Privacy Regulation", "GDPR, AI Act, enterprise data sovereignty requirements"),
        ("Voice + Multimodal", "CB Insights #1 prediction for 2026: multimodal agents win"),
        ("Local-First Demand", "On-prem AI is the #1 enterprise requirement for sensitive data"),
        ("Cost Optimization", "Companies seeking cheaper alternatives to GPT-4/Claude"),
    ]

    for j, (title, desc) in enumerate(drivers):
        y = Inches(2.7) + j * Inches(0.82)
        rect(slide, Inches(7.3), y, Inches(0.06), Inches(0.55), ACCENT)
        tb(slide, Inches(7.55), y, Inches(4.7), Inches(0.3),
           title, size=13, color=WHITE, bold=True)
        tb(slide, Inches(7.55), y + Inches(0.28), Inches(4.7), Inches(0.3),
           desc, size=10, color=GRAY)

    slide_number(slide, 7)


def slide_competitive(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    section_title(slide, "Competitive Landscape", "How Nexus compares")

    # Table
    headers = ["Feature", "Nexus", "ChatGPT", "Claude", "Cursor", "Copilot"]
    rows_data = [
        ["Self-evolving agent", "✓", "✗", "✗", "✗", "✗"],
        ["Local-first privacy", "✓", "✗", "✗", "✗", "✗"],
        ["Persistent memory & vault", "✓", "Limited", "Limited", "✗", "✗"],
        ["Skill self-authoring", "✓", "✗", "✗", "✗", "✗"],
        ["Voice + TTS", "✓", "App only", "✗", "✗", "✗"],
        ["Kanban / Project mgmt", "✓", "✗", "✗", "✗", "✗"],
        ["Model-agnostic", "✓", "✗", "✗", "✗", "✗"],
        ["Dream engine (idle learning)", "✓", "✗", "✗", "✗", "✗"],
        ["Offline capable", "✓", "✗", "✗", "Partial", "✗"],
        ["Cost-optimized routing", "✓", "✗", "✗", "✗", "✗"],
    ]

    cols_n = len(headers)
    rows_n = len(rows_data) + 1

    tbl_left = Inches(0.75)
    tbl_top = Inches(1.9)
    tbl_w = Inches(11.8)
    tbl_h = Inches(5.0)

    table_shape = slide.shapes.add_table(rows_n, cols_n, tbl_left, tbl_top, tbl_w, tbl_h)
    table = table_shape.table

    col_widths = [Inches(2.8), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Style header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT
            p.alignment = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_DIM if i == 1 else RGBColor(0x15, 0x23, 0x40)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Style data rows
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = FONT
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                if c == 1:
                    p.font.color.rgb = GREEN if val == "✓" else WHITE
                    p.font.bold = True
                elif val == "✓":
                    p.font.color.rgb = GREEN
                elif val == "✗":
                    p.font.color.rgb = RGBColor(0x47, 0x55, 0x69)
                else:
                    p.font.color.rgb = GOLD
            cell.fill.solid()
            if r % 2 == 0:
                cell.fill.fore_color.rgb = BG_CARD
            else:
                cell.fill.fore_color.rgb = RGBColor(0x0E, 0x19, 0x2E)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    slide_number(slide, 8)


def slide_moat(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    section_title(slide, "Competitive Moat", "Why this is defensible")

    moats = [
        ("Data Moat", "User vaults, knowledge graphs, dream journals — proprietary data that creates deep switching costs",
         ACCENT, 0.95),
        ("Self-Evolution IP", "Novel dream engine + autonomous skill authoring — no competitor has this",
         PURPLE, 0.88),
        ("Gateway Economics", "LiteLLM routing + GLM 5.1 = 60-80% cost advantage vs. OpenAI/Anthropic retail",
         GREEN, 0.82),
        ("Local-First Trust", "Enterprises won't send data to ChatGPT — Nexus keeps it local by design",
         GOLD, 0.78),
        ("Skills Ecosystem", "50 bundled + marketplace creates network effects: more skills → more users → more skills",
         CYAN, 0.72),
        ("Security Layer", "Prompt filtering, guard rails, secrets redaction — regulatory moat as AI governance tightens",
         RED_SOFT, 0.68),
    ]

    sy = Inches(1.9)
    bar_h = Inches(0.75)
    gap = Inches(0.12)
    max_w = Inches(10.5)
    sx = Inches(2.5)

    for i, (title, desc, color, pct) in enumerate(moats):
        y = sy + i * (bar_h + gap)

        # Title (left)
        tb(slide, Inches(0.75), y + Inches(0.05), Inches(1.7), Inches(0.35),
           title, size=12, color=color, bold=True, alignment=PP_ALIGN.RIGHT)

        # Bar background
        rect(slide, sx, y, max_w, bar_h, RGBColor(0x15, 0x23, 0x40))

        # Bar fill
        fill_w = int(max_w * pct)
        rect(slide, sx, y, fill_w, bar_h, color)

        # Description on bar
        tb(slide, sx + Inches(0.15), y + Inches(0.1), fill_w - Inches(0.3), Inches(0.55),
           desc, size=10, color=WHITE)

    slide_number(slide, 9)


def slide_business_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    section_title(slide, "Business Model", "Three-phase strategy for venture-scale growth")

    phases = [
        ("Phase 1", "Consumer SaaS", "Year 1-2", ACCENT,
         ["Free self-hosted + BYO model",
          "Pro $29/mo — Gateway + all skills",
          "Power $49/mo — Priority + sub-agents",
          "Revenue engine: Nexus model gateway"]),
        ("Phase 2", "Enterprise", "Year 2-3", PURPLE,
         ["Team $49/user/mo — Shared workspaces",
          "Business $99/user/mo — SSO + compliance",
          "Enterprise custom ($50K-$500K/yr)",
          "Custom skill development by Nexus team"]),
        ("Phase 3", "Platform", "Year 3+", GREEN,
         ["Skills marketplace (30% commission)",
          "Apps & recipes marketplace",
          "Prompt security product (standalone)",
          "Model marketplace for fine-tuned LLMs"]),
    ]

    pw = Inches(3.6)
    ph = Inches(4.6)
    gap = Inches(0.55)
    total = 3 * pw + 2 * gap
    sx = (SW - total) / 2
    sy = Inches(2.0)

    for i, (phase, title, period, color, bullets) in enumerate(phases):
        x = sx + i * (pw + gap)
        c = card(slide, x, sy, pw, ph, BG_CARD)

        # Phase label + top bar
        rect(slide, x, sy, pw, Inches(0.06), color)

        tb(slide, x + Inches(0.2), sy + Inches(0.2), Inches(2.0), Inches(0.3),
           phase, size=11, color=GRAY, bold=True)
        tb(slide, x + Inches(0.2), sy + Inches(0.5), Inches(3.0), Inches(0.4),
           title, size=20, color=WHITE, bold=True)
        tb(slide, x + Inches(0.2), sy + Inches(0.95), Inches(3.0), Inches(0.3),
           period, size=11, color=color, bold=True)

        accent_line(slide, x + Inches(0.2), sy + Inches(1.35), Inches(1.2), color, Pt(2))

        for j, bullet in enumerate(bullets):
            tb(slide, x + Inches(0.2), sy + Inches(1.55) + j * Inches(0.6),
               Inches(3.1), Inches(0.55),
               f"→  {bullet}", size=11, color=LIGHT)

        # Connector arrow
        if i < 2:
            ax = x + pw + Inches(0.1)
            ay = sy + ph / 2
            rect(slide, ax, ay, gap - Inches(0.2), Inches(0.05), GRAY)

    slide_number(slide, 10)


def slide_pricing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    section_title(slide, "Pricing", "Consumer tiers designed for conversion, enterprise for expansion")

    # Consumer tiers
    tiers = [
        ("Free", "$0", "/mo", GRAY, [
            "Self-hosted, BYO model",
            "Core agent features",
            "5 bundled skills",
            "Community support",
        ], False),
        ("Pro", "$29", "/mo", ACCENT, [
            "Nexus model gateway",
            "50+ skills library",
            "Cloud sync & mobile access",
            "Dream engine",
            "Voice (TTS + STT)",
        ], True),
        ("Power", "$49", "/mo", PURPLE, [
            "Everything in Pro",
            "Priority model routing",
            "Extended context windows",
            "Sub-agent spawning",
            "Local model management",
        ], False),
    ]

    tw = Inches(2.9)
    th = Inches(4.4)
    gap = Inches(0.45)
    total = 3 * tw + 2 * gap
    sx = (SW - total - Inches(2.5)) / 2
    sy = Inches(2.0)

    for i, (name, price, period, color, features, highlight) in enumerate(tiers):
        x = sx + i * (tw + gap)
        border = color if highlight else None
        c = card(slide, x, sy, tw, th, BG_CARD_ALT if highlight else BG_CARD, border_color=border)

        rect(slide, x, sy, tw, Inches(0.06), color)

        if highlight:
            tb(slide, x, sy + Inches(0.15), tw, Inches(0.25),
               "MOST POPULAR", size=9, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        tb(slide, x + Inches(0.2), sy + Inches(0.45), Inches(2.5), Inches(0.35),
           name, size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Price
        multi_tb(slide, x + Inches(0.2), sy + Inches(0.9), Inches(2.5), Inches(0.7), [
            (price, 36, color, True, PP_ALIGN.CENTER),
            (period, 12, GRAY, False, PP_ALIGN.CENTER),
        ])

        accent_line(slide, x + Inches(0.8), sy + Inches(1.7), Inches(1.3), RGBColor(0x1E, 0x29, 0x3B), Pt(1))

        for j, feat in enumerate(features):
            tb(slide, x + Inches(0.25), sy + Inches(1.9) + j * Inches(0.45),
               Inches(2.4), Inches(0.4),
               f"✓  {feat}", size=10, color=LIGHT)

    # Enterprise tiers (right side)
    ex = sx + 3 * (tw + gap) + Inches(0.3)
    ew = Inches(4.0)
    eh = Inches(4.4)

    ent_card = card(slide, ex, sy, ew, eh, BG_CARD, border_color=GOLD)
    rect(slide, ex, sy, ew, Inches(0.06), GOLD)

    tb(slide, ex + Inches(0.2), sy + Inches(0.2), ew - Inches(0.4), Inches(0.35),
       "ENTERPRISE", size=14, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    ent_tiers = [
        ("Team", "$49", "/user/mo", "Shared workspaces, team kanban"),
        ("Business", "$99", "/user/mo", "SSO, compliance, custom skills"),
        ("Enterprise", "Custom", "", "$50K-$500K/yr, on-prem, SLA"),
    ]

    for j, (name, price, period, desc) in enumerate(ent_tiers):
        ey = sy + Inches(0.65) + j * Inches(1.25)
        multi_tb(slide, ex + Inches(0.25), ey, ew - Inches(0.5), Inches(1.1), [
            (f"{name}  —  {price}{period}", 13, WHITE, True, PP_ALIGN.LEFT),
            (desc, 10, LIGHT, False, PP_ALIGN.LEFT),
        ])

    slide_number(slide, 11)


def slide_unit_economics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    section_title(slide, "Unit Economics", "Healthy margins driven by cost-optimized model routing")

    # Consumer Pro tier breakdown (left)
    tb(slide, Inches(0.75), Inches(1.9), Inches(5.5), Inches(0.35),
       "Consumer Pro Tier — $29/mo", size=16, color=WHITE, bold=True)

    econ_rows = [
        ("Revenue", "$29.00", WHITE, True),
        ("Inference COGS (GLM 5.1 wholesale)", "($10.00)", RED_SOFT, False),
        ("Infrastructure (sync, CDN, relay)", "($1.50)", RED_SOFT, False),
        ("Support & operations", "($1.00)", RED_SOFT, False),
        ("Total COGS", "($12.50)", RED_SOFT, True),
        ("", "", WHITE, False),
        ("Gross Profit", "$16.50/mo", GREEN, True),
        ("Gross Margin", "57%", GREEN, True),
    ]

    for j, (label, value, color, bold) in enumerate(econ_rows):
        y = Inches(2.4) + j * Inches(0.42)
        tb(slide, Inches(0.95), y, Inches(4.0), Inches(0.35),
           label, size=11, color=LIGHT if not bold else color, bold=bold)
        tb(slide, Inches(5.0), y, Inches(1.2), Inches(0.35),
           value, size=11, color=color, bold=bold, alignment=PP_ALIGN.RIGHT)

    # Right side: Key metrics
    tb(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(0.35),
       "Key SaaS Metrics (Targets)", size=16, color=WHITE, bold=True)

    metrics_card = card(slide, Inches(7.0), Inches(2.4), Inches(5.5), Inches(4.3), BG_CARD)

    metrics = [
        ("", "Year 1", "Year 2", "Year 3"),
        ("CAC (Consumer)", "$30-50", "$20-40", "$15-30"),
        ("LTV (24mo)", "$500-700", "$500-700", "$600-800"),
        ("LTV / CAC", "10-23x", "13-35x", "20-53x"),
        ("Monthly Churn", "5-8%", "3-5%", "2-3%"),
        ("Net Revenue Retention", "90%", "105%", "115%"),
        ("Gross Margin", "52%", "60%", "68%"),
    ]

    mc = len(metrics[0])
    mr = len(metrics)
    mt = slide.shapes.add_table(mr, mc, Inches(7.2), Inches(2.6), Inches(5.1), Inches(3.8))
    mtable = mt.table

    col_ws = [Inches(1.8), Inches(1.1), Inches(1.1), Inches(1.1)]
    for ci, w in enumerate(col_ws):
        mtable.columns[ci].width = w

    for ri, row in enumerate(metrics):
        for ci, val in enumerate(row):
            cell = mtable.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                if ri == 0:
                    p.font.bold = True
                    p.font.color.rgb = ACCENT_LIGHT
                elif ci == 0:
                    p.font.color.rgb = LIGHT
                else:
                    p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x0E, 0x19, 0x2E) if ri % 2 == 0 else BG_CARD
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    slide_number(slide, 12)


def slide_revenue(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    section_title(slide, "Revenue Projections", "Base scenario — path to $20M+ ARR")

    # Revenue table (left)
    headers = ["", "Year 1", "Year 2", "Year 3"]
    rows_data = [
        ["Pro subscribers", "1,000", "8,000", "35,000"],
        ["Power subscribers", "150", "1,500", "10,000"],
        ["Team seats", "—", "300", "3,000"],
        ["Enterprise contracts", "—", "5", "25"],
        ["", "", "", ""],
        ["Consumer MRR", "$33K", "$297K", "$1.42M"],
        ["Enterprise ARR", "—", "$500K", "$4.5M"],
        ["Skills marketplace", "—", "$20K", "$300K"],
        ["", "", "", ""],
        ["Total ARR", "$400K", "$4.1M", "$21.8M"],
        ["Gross Margin", "52%", "60%", "68%"],
        ["Gross Profit", "$208K", "$2.46M", "$14.8M"],
    ]

    rn = len(rows_data) + 1
    cn = len(headers)
    tbl = slide.shapes.add_table(rn, cn, Inches(0.75), Inches(1.9), Inches(6.3), Inches(5.0))
    table = tbl.table

    col_ws = [Inches(2.0), Inches(1.43), Inches(1.43), Inches(1.43)]
    for ci, w in enumerate(col_ws):
        table.columns[ci].width = w

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = ACCENT_LIGHT
            p.font.name = FONT
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x15, 0x23, 0x40)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for ri, row in enumerate(rows_data):
        is_total = row[0] in ("Total ARR", "Gross Profit")
        is_spacer = row[0] == ""
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT
                p.font.bold = is_total
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                if is_total:
                    p.font.color.rgb = GREEN
                elif ci == 0:
                    p.font.color.rgb = LIGHT
                else:
                    p.font.color.rgb = WHITE
            cell.fill.solid()
            if is_total:
                cell.fill.fore_color.rgb = RGBColor(0x13, 0x2A, 0x1C)
            elif ri % 2 == 0:
                cell.fill.fore_color.rgb = BG_CARD
            else:
                cell.fill.fore_color.rgb = RGBColor(0x0E, 0x19, 0x2E)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Revenue bar chart (right)
    chart_data = CategoryChartData()
    chart_data.categories = ["Year 1", "Year 2", "Year 3"]
    chart_data.add_series("Consumer ARR", (384, 3564, 17040))
    chart_data.add_series("Enterprise ARR", (0, 500, 4500))
    chart_data.add_series("Marketplace", (0, 20, 300))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, Inches(7.3), Inches(1.9), Inches(5.5), Inches(4.8),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(9)
    chart.legend.font.color.rgb = LIGHT
    chart.legend.include_in_layout = False

    # Style chart
    plot = chart.plots[0]
    plot.gap_width = 120

    series_colors = [ACCENT, PURPLE, GREEN]
    for i, series in enumerate(plot.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = series_colors[i]

    # Style axes
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(9)
    cat_axis.tick_labels.font.color.rgb = LIGHT
    cat_axis.format.line.fill.background()

    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(9)
    val_axis.tick_labels.font.color.rgb = GRAY
    val_axis.tick_labels.number_format = '$#,##0"K"'
    val_axis.format.line.fill.background()
    val_axis.has_major_gridlines = True
    val_axis.major_gridlines.format.line.color.rgb = RGBColor(0x1E, 0x29, 0x3B)

    # Chart title annotation
    tb(slide, Inches(7.3), Inches(6.75), Inches(5.5), Inches(0.3),
       "Annual Recurring Revenue ($K)", size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

    slide_number(slide, 13)


def slide_gtm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    section_title(slide, "Go-to-Market", "Community-led → Product-led → Enterprise")

    years = [
        ("Year 1", "Community-Led Growth", ACCENT,
         "Build the developer & power-user community",
         ["Open-source core (keep gateway proprietary)",
          "GitHub stars, Discord, conference talks",
          "\"Build your own AI agent\" tutorials",
          "Product Hunt + Hacker News launch",
          "Target: 10K free → 1K paying subscribers"]),
        ("Year 2", "Product-Led Growth", PURPLE,
         "Activate mobile & marketplace flywheels",
         ["Mobile app (iOS/Android via Capacitor)",
          "Skills marketplace launch",
          "Enterprise pilot program (5 design partners)",
          "Strategic partnerships (model providers)",
          "Target: 8K paying + 5 enterprise contracts"]),
        ("Year 3", "Enterprise Scale", GREEN,
         "Sales-led expansion with security product",
         ["Dedicated enterprise sales team",
          "SOC2 + GDPR compliance certifications",
          "Prompt security product launch",
          "Channel partnerships (SIs, cloud mktplaces)",
          "Target: 45K paying + 25 enterprise contracts"]),
    ]

    yw = Inches(3.6)
    yh = Inches(4.7)
    gap = Inches(0.55)
    total = 3 * yw + 2 * gap
    sx = (SW - total) / 2
    sy = Inches(1.9)

    for i, (year, title, color, subtitle, bullets) in enumerate(years):
        x = sx + i * (yw + gap)
        c = card(slide, x, sy, yw, yh, BG_CARD)

        # Top bar
        rect(slide, x, sy, yw, Inches(0.06), color)

        # Year badge
        badge = card(slide, x + Inches(0.2), sy + Inches(0.2), Inches(0.9), Inches(0.35), color)
        tb(slide, x + Inches(0.2), sy + Inches(0.2), Inches(0.9), Inches(0.35),
           year, size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        tb(slide, x + Inches(1.25), sy + Inches(0.2), Inches(2.1), Inches(0.35),
           title, size=14, color=WHITE, bold=True)

        tb(slide, x + Inches(0.2), sy + Inches(0.7), Inches(3.1), Inches(0.35),
           subtitle, size=10, color=GRAY)

        accent_line(slide, x + Inches(0.2), sy + Inches(1.15), Inches(1.0), color, Pt(2))

        for j, bullet in enumerate(bullets):
            tb(slide, x + Inches(0.2), sy + Inches(1.35) + j * Inches(0.58),
               Inches(3.1), Inches(0.5),
               f"→  {bullet}", size=10, color=LIGHT)

        # Connector
        if i < 2:
            ax = x + yw + Inches(0.1)
            ay = sy + yh / 2
            rect(slide, ax, ay, gap - Inches(0.2), Inches(0.05), GRAY)

    slide_number(slide, 14)


def slide_ask(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    section_title(slide, "The Ask", "Pre-seed round to launch and reach seed-ready metrics")

    # Big number
    tb(slide, Inches(0.75), Inches(1.9), Inches(4.5), Inches(1.0),
       "$1.5M", size=64, color=ACCENT, bold=True)
    tb(slide, Inches(0.75), Inches(2.85), Inches(4.5), Inches(0.4),
       "Pre-Seed  •  ~$12M post-money valuation", size=14, color=GRAY)

    # Use of funds (pie chart)
    chart_data = CategoryChartData()
    chart_data.categories = ["Engineering\n(3 hires)", "Infrastructure\n(Gateway, CDN)", "Go-to-Market\n(Community, content)", "Operations\n& Legal", "Runway\nReserve"]
    chart_data.add_series("Use of Funds", (750, 300, 225, 150, 75))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(0.5), Inches(3.4), Inches(5.0), Inches(3.8),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(9)
    chart.legend.font.color.rgb = LIGHT
    chart.legend.include_in_layout = False

    plot = chart.plots[0]
    series = plot.series[0]
    pie_colors = [ACCENT, PURPLE, GREEN, GOLD, GRAY]
    for idx, color in enumerate(pie_colors):
        point = series.points[idx]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color

    series.data_labels.font.size = Pt(9)
    series.data_labels.font.color.rgb = WHITE
    series.data_labels.number_format = '$#,##0"K"'

    # 18-month milestones (right side)
    mx = Inches(6.0)
    tb(slide, mx, Inches(1.9), Inches(6.5), Inches(0.4),
       "18-Month Milestones", size=18, color=WHITE, bold=True)

    milestones = [
        ("Public Launch", "Open-source core + Pro subscription live", ACCENT),
        ("10K+ Free Users", "Developer community with organic growth", CYAN),
        ("1K+ Paying Subscribers", "$400K ARR run-rate", GREEN),
        ("3+ Enterprise Pilots", "Design partner program validated", PURPLE),
        ("Seed-Ready Metrics", "Positioned for $8-15M Seed round", GOLD),
    ]

    for j, (title, desc, color) in enumerate(milestones):
        y = Inches(2.5) + j * Inches(0.95)
        mc = card(slide, mx, y, Inches(6.5), Inches(0.8), BG_CARD)
        rect(slide, mx, y, Inches(0.06), Inches(0.8), color)
        tb(slide, mx + Inches(0.25), y + Inches(0.08), Inches(5.8), Inches(0.3),
           title, size=13, color=WHITE, bold=True)
        tb(slide, mx + Inches(0.25), y + Inches(0.4), Inches(5.8), Inches(0.3),
           desc, size=10, color=GRAY)

    slide_number(slide, 15)


def slide_thankyou(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    rect(slide, Inches(0), Inches(0), SW, Inches(0.06), ACCENT)

    # Large NEXUS
    tb(slide, Inches(0.75), Inches(1.5), Inches(11.8), Inches(1.2),
       "NEXUS", size=80, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    accent_line(slide, Inches(5.5), Inches(2.8), Inches(2.3))

    tb(slide, Inches(1.5), Inches(3.1), Inches(10.3), Inches(0.7),
       "The Self-Evolving AI Agent", size=24, color=ACCENT_LIGHT,
       alignment=PP_ALIGN.CENTER)

    tb(slide, Inches(2.0), Inches(4.2), Inches(9.3), Inches(0.7),
       "Let's build the future of personal AI — together.", size=16, color=GRAY,
       alignment=PP_ALIGN.CENTER)

    # Contact placeholder
    contact_card = card(slide, Inches(3.5), Inches(5.3), Inches(6.3), Inches(1.2),
                        BG_CARD, border_color=ACCENT_DIM)
    multi_tb(slide, Inches(3.8), Inches(5.4), Inches(5.7), Inches(1.0), [
        ("Contact", 13, ACCENT_LIGHT, True, PP_ALIGN.CENTER),
        ("[founder@nexus-model.us]  •  [linkedin.com/in/founder]  •  [nexus-model.us]", 11, GRAY, False, PP_ALIGN.CENTER),
    ])

    slide_number(slide, 16)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    # Remove default slide layouts' background shapes
    # (we use blank layout = index 6)

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_product(prs)
    slide_architecture(prs)
    slide_dream(prs)
    slide_market(prs)
    slide_competitive(prs)
    slide_moat(prs)
    slide_business_model(prs)
    slide_pricing(prs)
    slide_unit_economics(prs)
    slide_revenue(prs)
    slide_gtm(prs)
    slide_ask(prs)
    slide_thankyou(prs)

    out = "/Users/nino/Code/nexus/Nexus_Investor_Deck.pptx"
    prs.save(out)
    print(f"✓ Deck saved to {out}")
    print(f"  {len(prs.slides)} slides generated")


if __name__ == "__main__":
    build()
